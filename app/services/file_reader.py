import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import openpyxl
import csv
import base64


def extract_text(file_path: str) -> str:
    
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}")
        return ""

    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".pdf":
            return _read_pdf(file_path)
        elif ext == ".docx":
            return _read_docx(file_path)
        elif ext == ".txt":
            return _read_txt(file_path)
        elif ext == ".pptx":
            return _read_pptx(file_path)
        elif ext == ".xlsx":
            return _read_xlsx(file_path)
        elif ext == ".csv":
            return _read_csv(file_path)
        elif ext == ".epub":
            return _read_epub(file_path)
        else:
            print(f"Unsupported file type: {ext}")
            return ""
    except Exception as e:
        print(f"Error reading {ext} file: {e}")
        return ""


def extract_images_from_pdf(file_path: str, max_images: int = 5) -> list:
    
    if not os.path.exists(file_path):
        return []

    images = []

    try:
        doc = fitz.open(file_path)

        for page_num in range(min(len(doc), max_images)):
            page = doc[page_num]
            page_text = page.get_text()

            # only render pages that likely contain figures
            has_figure = any(
                keyword in page_text.lower()
                for keyword in ["fig", "figure", "diagram", "chart", "table", "image"]
            )

            # always include first 3 pages, then only figure pages
            if page_num < 3 or has_figure:
                # render page at high resolution — 2x zoom for clarity
                mat = fitz.Matrix(2.0, 2.0)
                pix = page.get_pixmap(matrix=mat)

                # convert to PNG bytes then base64
                img_bytes = pix.tobytes("png")
                b64 = base64.b64encode(img_bytes).decode("utf-8")

                # find figure captions on this page
                captions = _find_all_captions(page_text)

                images.append({
                    "page": page_num + 1,
                    "index": page_num + 1,
                    "base64": b64,
                    "ext": "png",
                    "caption": ", ".join(captions) if captions else f"Page {page_num + 1}"
                })

                print(f"Rendered page {page_num + 1} as image ({len(img_bytes)} bytes)")

        doc.close()
        print(f"Rendered {len(images)} pages as images")

    except Exception as e:
        print(f"Error rendering PDF pages: {e}")

    return images


def _find_all_captions(page_text: str) -> list:
    
    captions = []
    lines = page_text.split("\n")
    for line in lines:
        line_lower = line.lower().strip()
        if line_lower.startswith("fig") or "figure" in line_lower[:10]:
            captions.append(line.strip()[:100])
    return captions


def _read_pdf(file_path: str) -> str:
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text() + "\n"
    doc.close()
    return _clean(text)


def _read_docx(file_path: str) -> str:
    doc = Document(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return _clean("\n".join(paragraphs))


def _read_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return _clean(f.read())


def _read_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text = ""
    for slide_num, slide in enumerate(prs.slides, start=1):
        text += f"\n--- Slide {slide_num} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + "\n"
    return _clean(text)


def _read_xlsx(file_path: str) -> str:
    wb = openpyxl.load_workbook(file_path, data_only=True)
    text = ""
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        text += f"\n--- Sheet: {sheet_name} ---\n"
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
            if row_text.strip():
                text += row_text + "\n"
    return _clean(text)


def _read_csv(file_path: str) -> str:
    text = ""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            text += "\t".join(row) + "\n"
    return _clean(text)


def _read_epub(file_path: str) -> str:
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup

        book = epub.read_epub(file_path)
        text = ""
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), "html.parser")
                text += soup.get_text() + "\n"
        return _clean(text)
    except ImportError:
        print("ebooklib not installed. Run: pip install ebooklib")
        return ""


def _clean(text: str) -> str:
    lines = [line.strip() for line in text.split("\n") if line.strip()]
    return "\n".join(lines)


def extract_text_from_pdf(file_path: str) -> str:
    return _read_pdf(file_path)


SUPPORTED_EXTENSIONS = [".pdf", ".docx", ".txt", ".pptx", ".xlsx", ".csv", ".epub"]


def is_supported(filename: str) -> bool:
    ext = os.path.splitext(filename)[1].lower()
    return ext in SUPPORTED_EXTENSIONS
